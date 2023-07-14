from PyPDF2 import PdfReader
from docx import Document
from docx2pdf import convert
import tabula
import pandas as pd
from openpyxl import load_workbook
from reportlab.platypus import SimpleDocTemplate, Table
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from pdf2image import convert_from_path
import comtypes.client


# Button 1 : PDF to Word format


def convert_pdf_to_word(pdf_path, docx_path):
    try:
        pdf_reader = PdfReader(pdf_path)

        document = Document()
        for page in range(len(pdf_reader.pages)):
            page_content = pdf_reader.pages[page].extract_text()
            document.add_paragraph(page_content)

        document.save(docx_path)

        return True
    except Exception as e:
        print(f"Error converting PDF to Word: {str(e)}")
        return False


# Button 2 : WORD to PDF format

def convert_word_to_pdf(docx_path, pdf_path):
    try:
        convert(docx_path, pdf_path)
        return True
    except Exception as e:
        print(f"Error converting Word to PDF: {str(e)}")
        return False


# Button 3 : EXCEL to PDF format

def convert_excel_to_pdf(excel_path, pdf_path):
    try:
        # Load the Excel workbook
        workbook = load_workbook(excel_path)

        # Select the active sheet
        sheet = workbook.active

        # Get the data from the sheet
        data = sheet.values

        # Create a list of rows from the sheet data
        rows = list(data)

        # Create the PDF document
        doc = SimpleDocTemplate(pdf_path, pagesize=letter)

        # Create the table from the rows
        table = Table(rows)

        # Set the table style (optional)
        table.setStyle([('GRID', (0, 0), (-1, -1), 1, (0.5, 0.5, 0.5))])

        # Build the PDF document with the table
        elements = [table]
        doc.build(elements)

        return True
    except Exception as e:
        print(f"Error converting Excel to PDF: {str(e)}")
        return False


# Button 4 : PDF to EXCEL format


def convert_pdf_to_excel(pdf_path, excel_path):
    try:
        # Read the PDF and extract tables
        tables = tabula.read_pdf(pdf_path, pages='all')

        # Convert each table to an Excel sheet
        with pd.ExcelWriter(excel_path) as writer:
            for i, table in enumerate(tables, start=1):
                table.to_excel(writer, sheet_name=f"Sheet{i}", index=False)

        return True
    except Exception as e:
        print(f"Error converting PDF to Excel: {str(e)}")
        return False


# Button 5 : PDF to PowerPoint format

def convert_pdf_to_ppt(pdf_path, pptx_path):
    try:
        # Convert PDF pages to images
        images = convert_from_path(pdf_path)

        # Create a new PowerPoint presentation
        presentation = Presentation()

        # Add each image as a new slide to the presentation
        for image in images:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.add_picture(image, 0, 0)

        # Save the PowerPoint presentation
        presentation.save(pptx_path)

        return True
    except Exception as e:
        print(f"Error converting PDF to PowerPoint: {str(e)}")
        return False


# button 6 : Powerpoint to PDF format

def convert_ppt_to_pdf(pptx_path, pdf_path):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(pptx_path)
    presentation.ExportAsFixedFormat(pdf_path, 2)  # 2 represents the PDF format

    presentation.Close()
    powerpoint.Quit()

    return True
